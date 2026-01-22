from . import Document, Page, DocumentLoader
from pptx import Presentation
from .image_utils import image_to_base64_uri
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .vision_utils import process_image_with_vision

class PptxLoader(DocumentLoader):
    """A loader for Microsoft PowerPoint (.pptx) files."""

    def load(self, source: str) -> Document:
        """
        Reads text and extracts images from a .pptx file on a slide-by-slide basis.
        
        Returns Document with structured_elements format using Vision LLM.
        """
        try:
            prs = Presentation(source)
            doc_pages = []

            for i, slide in enumerate(prs.slides):
                structured_elements = []

                for shape in slide.shapes:
                    # Extract text elements
                    if hasattr(shape, "text") and shape.text.strip():
                        structured_elements.append({
                            "type": "text",
                            "content": shape.text
                        })
                    
                    # Extract image elements
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                            base64_image_uri = image_to_base64_uri(image_bytes)
                            
                            if base64_image_uri:
                                vision_result = process_image_with_vision(image_bytes)
                                
                                structured_elements.append({
                                    "type": "image",
                                    "base64": base64_image_uri,
                                    "vision_description": vision_result["description"],
                                    "vision_tokens": vision_result["total_tokens"],
                                    "vision_cost": vision_result["cost"]
                                })
                        except Exception as e:
                            print(f"Warning: Could not process an image on slide {i+1}: {e}")

                doc_pages.append(Page(
                    page_number=i + 1,
                    structured_elements=structured_elements
                ))
            
            print(f"✅ Successfully loaded {len(doc_pages)} slides from PPTX")
            return Document(source=source, pages=doc_pages)
            
        except Exception as e:
            print(f"Error reading PPTX file: {str(e)}")
            raise e
